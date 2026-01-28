"""
Script to generate College Management System PowerPoint Presentation
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Add background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 51, 102)
    shape.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(200, 200, 200)
    p2.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(0, 51, 102)
    header.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Content
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(20)
        p.space_after = Pt(12)

def add_table_slide(prs, title, headers, rows):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(0, 51, 102)
    header.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Table
    cols = len(headers)
    table = slide.shapes.add_table(len(rows) + 1, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(0.5 * (len(rows) + 1))).table

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)

    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, val in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs, "College Management System", "Software Requirements Specification")

    # Slide 2: Introduction
    add_content_slide(prs, "Introduction", [
        "Purpose: Automate and streamline academic and administrative processes",
        "Designed for educational institutions",
        "Comprehensive management of students, staff, and operations",
        "Web-based application with role-based access control"
    ])

    # Slide 3: Technology Stack
    add_content_slide(prs, "Technology Stack", [
        "Backend: Python Flask 3.0",
        "Database: SQLite with SQLAlchemy ORM",
        "Authentication: Flask-Login",
        "Forms: Flask-WTF, WTForms",
        "Frontend: HTML Templates (Jinja2), CSS, JavaScript",
        "Security: Werkzeug password hashing, CSRF protection"
    ])

    # Slide 4: User Roles
    add_table_slide(prs, "User Roles",
        ["Role", "Description"],
        [
            ["Student", "View attendance, marks, fees, timetable; submit complaints"],
            ["Staff", "Mark attendance, enter marks, manage library"],
            ["HOD", "Staff privileges + department-level management"],
            ["Management", "Full administrative access to all modules"]
        ])

    # Slide 5: Authentication
    add_content_slide(prs, "Authentication Methods", [
        "Students: Login using Roll Number + Date of Birth",
        "Staff/HOD/Management: Username + Password + Role Selection",
        "Password hashing using Werkzeug security",
        "Session cookie security (HTTPOnly, SameSite)",
        "Role-based access control on all routes"
    ])

    # Slide 6: Core Modules Overview
    add_content_slide(prs, "Core Modules", [
        "User Management - Account creation and role assignment",
        "Department Management - Departments and HOD assignment",
        "Subject Management - Courses with staff assignments",
        "Attendance Module - Session-based attendance tracking",
        "Marks/Examination Module - Exam results management",
        "Fees Module - Fee structures and payment tracking"
    ])

    # Slide 7: Additional Modules
    add_content_slide(prs, "Additional Modules", [
        "Library Module - Book inventory and borrowing system",
        "Complaints Module - Student complaint submission and tracking",
        "Feedback Module - Staff feedback from students",
        "Notice Board Module - Announcements and notices",
        "Timetable Module - Class schedules management",
        "Notifications Module - System-generated alerts"
    ])

    # Slide 8: Attendance Module
    add_content_slide(prs, "Attendance Module", [
        "Staff Functions:",
        "  - Create attendance sessions for subjects",
        "  - Mark attendance (Present/Absent)",
        "  - View attendance records",
        "Student Functions:",
        "  - View personal attendance records",
        "  - View attendance summary/percentage"
    ])

    # Slide 9: Marks Module
    add_content_slide(prs, "Marks/Examination Module", [
        "Staff Functions:",
        "  - Create exams (Internal, Mid-term, Final)",
        "  - Enter marks for students",
        "  - Generate mark reports",
        "Student Functions:",
        "  - View personal exam results",
        "  - View semester-wise results"
    ])

    # Slide 10: Fees Module
    add_content_slide(prs, "Fees Module", [
        "Management Functions:",
        "  - Define fee structures (tuition, exam fees)",
        "  - Set due dates and academic year",
        "  - Upload fee payment records",
        "  - Generate pending fee reports",
        "Student Functions:",
        "  - View fee details and payment status",
        "  - View payment history"
    ])

    # Slide 11: Library Module
    add_content_slide(prs, "Library Module", [
        "Staff/Management Functions:",
        "  - Add, update, delete books",
        "  - Manage book inventory",
        "  - Issue books and track returns",
        "Student Functions:",
        "  - Search books by title, author, ISBN",
        "  - View book availability",
        "  - View personal borrowed books"
    ])

    # Slide 12: Database Schema
    add_table_slide(prs, "Database Schema - Core Tables",
        ["Table", "Description"],
        [
            ["users", "Base user accounts with role"],
            ["students", "Student profiles linked to users"],
            ["staff", "Staff profiles linked to users"],
            ["departments", "Department information"],
            ["subjects", "Subject/course information"],
            ["staff_subjects", "Staff-subject assignments (M:N)"]
        ])

    # Slide 13: Database Schema - Module Tables
    add_table_slide(prs, "Database Schema - Module Tables",
        ["Table", "Description"],
        [
            ["attendance_sessions", "Attendance session records"],
            ["attendance_records", "Individual attendance entries"],
            ["exams / marks", "Examination and student marks"],
            ["fee_structures / student_fees", "Fee definitions and payments"],
            ["books / book_issues", "Library inventory and borrowing"],
            ["complaints / feedbacks", "Student complaints and feedback"],
            ["notices / timetables", "Announcements and schedules"]
        ])

    # Slide 14: System Architecture
    add_content_slide(prs, "System Architecture", [
        "App Factory Pattern with Flask Blueprints",
        "Modular structure: app/, blueprints/, models/, templates/",
        "Each module has routes.py and forms.py",
        "Centralized extensions (SQLAlchemy, Login, CSRF)",
        "Role-based decorators for access control",
        "Business logic in services layer"
    ])

    # Slide 15: API Endpoints
    add_table_slide(prs, "API Endpoints Summary",
        ["Module", "Prefix", "Key Endpoints"],
        [
            ["Auth", "/", "/login, /logout"],
            ["Attendance", "/attendance", "/mark, /my-attendance"],
            ["Marks", "/marks", "/create-exam, /enter, /my-results"],
            ["Fees", "/fees", "/upload, /pending-report, /my-fees"],
            ["Library", "/library", "/search, /manage, /issue"],
            ["Complaints", "/complaints", "/submit, /my-complaints"]
        ])

    # Slide 16: Security Features
    add_content_slide(prs, "Security Features", [
        "Password hashing using Werkzeug security",
        "CSRF protection on all forms",
        "Session cookie security (HTTPOnly, SameSite)",
        "Role-based access control on all routes",
        "Input validation with WTForms",
        "SQLAlchemy ORM prevents SQL injection"
    ])

    # Slide 17: Non-Functional Requirements
    add_content_slide(prs, "Non-Functional Requirements", [
        "Security: Password hashing, CSRF, session security",
        "Usability: Responsive interface, role-specific dashboards",
        "Performance: Lazy loading, database indexing",
        "Maintainability: Modular architecture, separation of concerns",
        "Scalability: Blueprint-based structure for easy extension"
    ])

    # Slide 18: Thank You
    add_title_slide(prs, "Thank You", "College Management System - January 2026")

    prs.save('College_Management_System.pptx')
    print("Presentation saved as 'College_Management_System.pptx'")

if __name__ == "__main__":
    create_presentation()
